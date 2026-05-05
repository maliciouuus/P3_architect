from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BG   = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT    = RGBColor(0x89, 0xB4, 0xFA)
RED       = RGBColor(0xF3, 0x8B, 0xA8)
GREEN     = RGBColor(0xA6, 0xE3, 0xA1)
YELLOW    = RGBColor(0xF9, 0xE2, 0xAF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xBA, 0xC2, 0xDE)
ORANGE    = RGBColor(0xFA, 0xB3, 0x87)

blank_layout = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG
    return slide


def text_box(slide, text, x, y, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def rect(slide, x, y, w, h, fill_color, radius=False):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_bullet(slide, items, x, y, w, h, size=16, color=WHITE, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = "▸  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color


# ── SLIDE 1 : Titre ─────────────────────────────────────────────────────────
slide = add_slide()
rect(slide, 0, 2.8, 13.33, 0.06, ACCENT)
text_box(slide, "Renote — Analyse d'architecture", 0.5, 1.2, 12, 1.2,
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text_box(slide, "Problèmes actuels · Migration SOLID + MVC · API REST",
         0.5, 2.5, 12, 0.7, size=20, color=ACCENT, align=PP_ALIGN.CENTER)
text_box(slide, "Projet P3 — OpenClassrooms", 0.5, 6.5, 12, 0.6,
         size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 2 : C'est quoi l'appli ? ──────────────────────────────────────────
slide = add_slide()
text_box(slide, "L'application Renote", 0.5, 0.3, 12, 0.8, size=30, bold=True, color=ACCENT)
rect(slide, 0.5, 1.1, 12.3, 0.04, ACCENT)

text_box(slide, "Une app de prise de notes avec tags — stack Laravel + Livewire",
         0.5, 1.3, 12, 0.6, size=18, color=GRAY)

boxes = [
    ("👤 User",     "Inscription / Connexion",         0.4,  2.2),
    ("📝 Notes",    "Créer, lire, supprimer des notes", 3.4,  2.2),
    ("🏷️ Tags",     "Créer des tags et les associer",   6.4,  2.2),
    ("🗄️ SQLite",   "Base de données embarquée",        9.4,  2.2),
]
for title, desc, bx, by in boxes:
    rect(slide, bx, by, 2.7, 1.6, RGBColor(0x31, 0x32, 0x44))
    text_box(slide, title, bx+0.1, by+0.1, 2.5, 0.5, size=16, bold=True, color=ACCENT)
    text_box(slide, desc,  bx+0.1, by+0.6, 2.5, 0.8, size=13, color=WHITE)

text_box(slide, "Architecture actuelle : server-driven (tout en PHP, HTML généré côté serveur)",
         0.5, 4.1, 12, 0.5, size=16, color=YELLOW)

text_box(slide, "Fichiers clés :", 0.5, 4.8, 12, 0.4, size=16, bold=True, color=WHITE)
add_bullet(slide, [
    "app/Livewire/Notes.php  —  composant principal",
    "app/Livewire/TagForm.php  —  formulaire de tag",
    "app/Models/Note.php + Tag.php  —  modèles Eloquent",
    "routes/web.php  —  routes HTML uniquement (pas d'API)",
], 0.5, 5.2, 12, 2.0, size=15, color=GRAY)


# ── SLIDE 3 : Problème 1 — Single Responsibility ────────────────────────────
slide = add_slide()
text_box(slide, "Problème #1 — Violation du principe S (Single Responsibility)",
         0.5, 0.3, 12, 0.8, size=26, bold=True, color=RED)
rect(slide, 0.5, 1.1, 12.3, 0.04, RED)

text_box(slide, "Notes.php fait TOUT en même temps :", 0.5, 1.3, 12, 0.5,
         size=18, bold=True, color=WHITE)

responsibilities = [
    ("Validation",      "Règles $rules + appel à validate()",                   0.4, 2.0, YELLOW),
    ("Accès BDD",       "Note::create(), Note::where(), Tag::all()",             3.6, 2.0, RED),
    ("État de l'UI",    "$this->text, $this->tag_id, $this->notes",             6.8, 2.0, ORANGE),
    ("Messages flash",  "session()->flash('message', ...)",                     0.4, 4.0, ACCENT),
    ("Chargement data", "loadNotes(), refreshTags()",                           3.6, 4.0, GREEN),
]
for title, desc, bx, by, col in responsibilities:
    rect(slide, bx, by, 2.8, 1.7, RGBColor(0x31, 0x32, 0x44))
    rect(slide, bx, by, 2.8, 0.35, col)
    text_box(slide, title, bx+0.1, by+0.05, 2.6, 0.3, size=13, bold=True, color=DARK_BG)
    text_box(slide, desc,  bx+0.1, by+0.45, 2.6, 1.1, size=12, color=GRAY)

text_box(slide,
         "⚠  Une seule classe = 5 responsabilités. Modifier la validation oblige à toucher au même fichier que la BDD.",
         0.5, 6.0, 12, 0.8, size=14, color=YELLOW)


# ── SLIDE 4 : Problème 2 — Pas de Service Layer ─────────────────────────────
slide = add_slide()
text_box(slide, "Problème #2 — Logique métier dans le composant UI",
         0.5, 0.3, 12, 0.8, size=26, bold=True, color=RED)
rect(slide, 0.5, 1.1, 12.3, 0.04, RED)

text_box(slide, "Actuellement :", 0.5, 1.3, 5, 0.4, size=18, bold=True, color=WHITE)
rect(slide, 0.5, 1.8, 5.5, 3.2, RGBColor(0x31, 0x32, 0x44))
text_box(slide, "Notes.php (Livewire)\n\n  save()\n    → validate()\n    → Note::create()\n    → session()->flash()\n    → loadNotes()\n\n  delete()\n    → Note::where()->delete()\n    → loadNotes()",
         0.6, 1.9, 5.3, 3.0, size=14, color=ORANGE)

text_box(slide, "Conséquences :", 6.5, 1.3, 6, 0.4, size=18, bold=True, color=WHITE)
add_bullet(slide, [
    "Impossible de réutiliser 'créer une note' depuis une API",
    "Impossible de tester la logique métier sans UI",
    "Si on change le framework, tout est à réécrire",
    "Duplications inévitables si on ajoute une 2ème interface",
], 6.5, 1.8, 6.3, 3.5, size=15, color=GRAY)

text_box(slide, "Solution → créer un NoteService et TagService qui contiennent la logique, injectés dans les contrôleurs.",
         0.5, 5.5, 12, 0.8, size=15, color=GREEN)


# ── SLIDE 5 : Problème 3 — Pas d'API ────────────────────────────────────────
slide = add_slide()
text_box(slide, "Problème #3 — Aucune API REST",
         0.5, 0.3, 12, 0.8, size=26, bold=True, color=RED)
rect(slide, 0.5, 1.1, 12.3, 0.04, RED)

text_box(slide, "routes/web.php — Routes disponibles actuellement :",
         0.5, 1.3, 12, 0.5, size=16, bold=True, color=WHITE)

routes_current = [
    ("GET", "/",                  "Page d'accueil"),
    ("GET", "/dashboard",         "Dashboard (HTML Blade)"),
    ("GET", "/notes",             "Redirige vers dashboard"),
    ("GET", "/tags",              "Redirige vers dashboard"),
    ("GET", "/settings/profile",  "Profil utilisateur"),
]
for i, (method, path, desc) in enumerate(routes_current):
    y = 2.0 + i * 0.55
    rect(slide, 0.5, y, 1.2, 0.42, RGBColor(0x45, 0x47, 0x5A))
    text_box(slide, method, 0.55, y+0.05, 1.1, 0.35, size=13, bold=True, color=YELLOW)
    text_box(slide, path,   1.8,  y+0.05, 3.5, 0.35, size=13, color=ACCENT)
    text_box(slide, desc,   5.4,  y+0.05, 7.0, 0.35, size=13, color=GRAY)

text_box(slide, "❌  Zéro endpoint API. Le front React ne peut rien appeler.",
         0.5, 5.0, 12, 0.5, size=16, color=RED)

text_box(slide, "Endpoints à créer :", 0.5, 5.6, 12, 0.4, size=16, bold=True, color=GREEN)
add_bullet(slide, [
    "GET /api/notes  ·  POST /api/notes  ·  DELETE /api/notes/{id}",
    "GET /api/tags   ·  POST /api/tags",
], 0.5, 6.1, 12, 0.9, size=15, color=GREEN)


# ── SLIDE 6 : Problème 4 — Couplage Livewire ────────────────────────────────
slide = add_slide()
text_box(slide, "Problème #4 — Front couplé au back (Livewire)",
         0.5, 0.3, 12, 0.8, size=26, bold=True, color=RED)
rect(slide, 0.5, 1.1, 12.3, 0.04, RED)

text_box(slide, "notes.blade.php — Vue actuelle :", 0.5, 1.3, 12, 0.5, size=16, bold=True, color=WHITE)

rect(slide, 0.5, 1.8, 7, 3.5, RGBColor(0x31, 0x32, 0x44))
text_box(slide,
         '<form wire:submit.prevent="save">\n'
         '  <textarea wire:model="text"></textarea>\n'
         '  <select wire:model="tag_id">...</select>\n'
         '  <button type="submit">Add Note</button>\n'
         '</form>\n\n'
         '@foreach ($notes as $note)\n'
         '  <button wire:click="delete({{$note->id}})">\n'
         '    Delete\n'
         '  </button>\n'
         '@endforeach',
         0.6, 1.9, 6.8, 3.3, size=13, color=ORANGE)

add_bullet(slide, [
    "wire:submit, wire:model, wire:click → 100% dépendant de Livewire",
    "HTML généré côté serveur → impossible à utiliser avec React",
    "Pas de séparation View / ViewModel",
    "Impossible d'avoir une app mobile ou un 2ème client",
    "Tout changement UI implique de modifier du PHP",
], 7.8, 1.8, 5.0, 4.0, size=15, color=GRAY)

text_box(slide, "Solution → remplacer les vues Blade/Livewire par des composants React qui appellent l'API REST.",
         0.5, 5.8, 12, 0.8, size=15, color=GREEN)


# ── SLIDE 7 : Récap STUPID ──────────────────────────────────────────────────
slide = add_slide()
text_box(slide, "Récap — Les antipatterns STUPID détectés",
         0.5, 0.3, 12, 0.8, size=28, bold=True, color=YELLOW)
rect(slide, 0.5, 1.1, 12.3, 0.04, YELLOW)

stupid = [
    ("S — Singleton",          "Non applicable ici",                                    GREEN,  False),
    ("T — Tight Coupling",     "Livewire couple UI + logique + BDD dans un seul fichier", RED,  True),
    ("U — Untestability",      "Impossible de tester Notes.php sans lancer tout Laravel", RED,  True),
    ("P — Premature Optim.",   "Non applicable ici",                                    GREEN,  False),
    ("I — Indescriptive Names","Noms corrects (Note, Tag, save, delete…)",              GREEN,  False),
    ("D — Duplication",        "Tag::all() dupliqué dans mount() et refreshTags()",     ORANGE, True),
]
for i, (letter, desc, col, problem) in enumerate(stupid):
    y = 1.4 + i * 0.83
    rect(slide, 0.5, y, 12.3, 0.75, RGBColor(0x31, 0x32, 0x44))
    rect(slide, 0.5, y, 0.08, 0.75, col)
    icon = "❌" if problem else "✅"
    text_box(slide, icon + "  " + letter, 0.7, y+0.1, 3.5, 0.55, size=15, bold=True, color=col)
    text_box(slide, desc, 4.3, y+0.1, 8.5, 0.55, size=14, color=WHITE if problem else GRAY)


# ── SLIDE 8 : Architecture cible ────────────────────────────────────────────
slide = add_slide()
text_box(slide, "Architecture cible — MVC + Services + API REST",
         0.5, 0.3, 12, 0.8, size=26, bold=True, color=GREEN)
rect(slide, 0.5, 1.1, 12.3, 0.04, GREEN)

layers = [
    ("FRONT (React)",           "Composants React + Zustand",  "routes séparées",          ACCENT,  0.4),
    ("API Layer",               "routes/api.php",              "GET/POST/PUT/DELETE JSON",  YELLOW,  2.0),
    ("Controllers",             "NoteController, TagController","Reçoit la requête HTTP",   ORANGE,  3.3),
    ("Services",                "NoteService, TagService",     "Logique métier pure",       GREEN,   4.6),
    ("Models / ORM",            "Note, Tag, User",             "Eloquent — accès BDD",      ACCENT,  5.9),
]
for name, files, role, col, y in layers:
    rect(slide, 0.4, y, 12.5, 1.0, RGBColor(0x31, 0x32, 0x44))
    rect(slide, 0.4, y, 0.1,  1.0, col)
    text_box(slide, name,  0.65, y+0.1, 3.0, 0.35, size=15, bold=True, color=col)
    text_box(slide, files, 0.65, y+0.5, 3.0, 0.4,  size=12, color=GRAY)
    text_box(slide, role,  4.0,  y+0.25,8.5, 0.5,  size=14, color=WHITE)

    if y < 5.9:
        text_box(slide, "↓", 6.2, y+1.0, 1, 0.3, size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 9 : Plan d'action ──────────────────────────────────────────────────
slide = add_slide()
text_box(slide, "Plan d'action — Ce qu'on va faire",
         0.5, 0.3, 12, 0.8, size=28, bold=True, color=ACCENT)
rect(slide, 0.5, 1.1, 12.3, 0.04, ACCENT)

steps = [
    ("1", "Créer NoteController + TagController",  "app/Http/Controllers/",   ACCENT),
    ("2", "Créer NoteService + TagService",         "app/Services/",           GREEN),
    ("3", "Définir les routes API REST",            "routes/api.php",          YELLOW),
    ("4", "Tester les endpoints avec Postman",      "GET, POST, DELETE",        ORANGE),
    ("5", "Créer le front React",                  "resources/js/",            ACCENT),
    ("6", "Connecter React à l'API",               "fetch / axios",            GREEN),
    ("7", "Supprimer les composants Livewire",      "app/Livewire/ → obsolète", RED),
]
for step, title, detail, col in steps:
    y = 1.3 + (int(step)-1) * 0.75
    rect(slide, 0.4, y, 0.55, 0.6, col)
    text_box(slide, step, 0.45, y+0.05, 0.45, 0.5, size=18, bold=True,
             color=DARK_BG, align=PP_ALIGN.CENTER)
    text_box(slide, title,  1.1, y+0.05, 7.5, 0.5, size=15, bold=True, color=WHITE)
    text_box(slide, detail, 8.8, y+0.05, 4.0, 0.5, size=13, color=GRAY)


# ── SLIDE 10 : NoteService — code annoté ─────────────────────────────────────
slide = add_slide()
text_box(slide, "Nouveau fichier — NoteService.php", 0.5, 0.3, 12, 0.8, size=26, bold=True, color=GREEN)
rect(slide, 0.5, 1.1, 12.3, 0.04, GREEN)

text_box(slide, "app/Services/NoteService.php", 0.5, 1.3, 12, 0.4, size=14, color=GRAY)

rect(slide, 0.5, 1.8, 7.8, 5.0, RGBColor(0x1E, 0x1E, 0x2E))
code_lines = (
    "class NoteService\n"
    "{\n"
    "  // Récupère toutes les notes de\n"
    "  // l'utilisateur connecté\n"
    "  public function getAllForUser()\n"
    "  {\n"
    "    return Note::with('tag')\n"
    "      ->where('user_id', Auth::id())\n"
    "      ->latest()->get();\n"
    "  }\n\n"
    "  // Crée une note — user_id injecté ici\n"
    "  public function create(array $data)\n"
    "  {\n"
    "    return Note::create([\n"
    "      'user_id' => Auth::id(),\n"
    "      'tag_id'  => $data['tag_id'],\n"
    "      'text'    => $data['text'],\n"
    "    ]);\n"
    "  }\n\n"
    "  // Supprime — vérifie l'appartenance\n"
    "  public function delete(int $noteId): bool\n"
    "  {\n"
    "    return (bool) Note::where('id', $noteId)\n"
    "      ->where('user_id', Auth::id())\n"
    "      ->delete();\n"
    "  }\n"
    "}"
)
text_box(slide, code_lines, 0.65, 1.9, 7.5, 4.8, size=11, color=ORANGE)

annotations = [
    (1.95, "✅ S — une seule responsabilité : la logique métier des notes"),
    (2.75, "✅ D — Auth::id() vient du framework, pas hardcodé"),
    (3.95, "✅ Réutilisable depuis n'importe quel contrôleur ou test"),
    (5.35, "✅ Sécurité : vérifie le propriétaire DANS le service"),
]
for ay, txt in annotations:
    rect(slide, 8.5, ay, 4.7, 0.55, RGBColor(0x31, 0x32, 0x44))
    text_box(slide, txt, 8.6, ay+0.05, 4.5, 0.45, size=12, color=GREEN)


# ── SLIDE 11 : NoteController — code annoté ──────────────────────────────────
slide = add_slide()
text_box(slide, "Nouveau fichier — NoteController.php", 0.5, 0.3, 12, 0.8, size=26, bold=True, color=ORANGE)
rect(slide, 0.5, 1.1, 12.3, 0.04, ORANGE)

text_box(slide, "app/Http/Controllers/Api/NoteController.php", 0.5, 1.3, 12, 0.4, size=14, color=GRAY)

rect(slide, 0.5, 1.8, 7.8, 5.0, RGBColor(0x1E, 0x1E, 0x2E))
code_lines = (
    "class NoteController extends Controller\n"
    "{\n"
    "  // NoteService injecté automatiquement\n"
    "  // par Laravel (principe D)\n"
    "  public function __construct(\n"
    "    private NoteService $noteService\n"
    "  ) {}\n\n"
    "  // GET /api/notes\n"
    "  public function index(): JsonResponse\n"
    "  {\n"
    "    $notes = $this->noteService\n"
    "               ->getAllForUser();\n"
    "    return response()->json([\n"
    "      'status' => 'success',\n"
    "      'data'   => $notes,\n"
    "    ]);\n"
    "  }\n\n"
    "  // POST /api/notes\n"
    "  public function store(Request $req)\n"
    "  {\n"
    "    $validated = $req->validate([...]);\n"
    "    $note = $this->noteService\n"
    "              ->create($validated);\n"
    "    return response()->json([...], 201);\n"
    "  }\n"
    "}"
)
text_box(slide, code_lines, 0.65, 1.9, 7.5, 4.8, size=11, color=ORANGE)

annotations = [
    (2.1,  "✅ D — pas de new NoteService() : injection auto"),
    (3.3,  "✅ S — le contrôleur ne fait que recevoir / répondre"),
    (4.3,  "✅ Réponse JSON uniforme : { status, data }"),
    (5.5,  "✅ Validation à la frontière (contrôleur), pas dans le service"),
]
for ay, txt in annotations:
    rect(slide, 8.5, ay, 4.7, 0.55, RGBColor(0x31, 0x32, 0x44))
    text_box(slide, txt, 8.6, ay+0.05, 4.5, 0.45, size=12, color=ORANGE)


# ── SLIDE 12 : routes/api.php ─────────────────────────────────────────────────
slide = add_slide()
text_box(slide, "Nouveau fichier — routes/api.php", 0.5, 0.3, 12, 0.8, size=26, bold=True, color=YELLOW)
rect(slide, 0.5, 1.1, 12.3, 0.04, YELLOW)

text_box(slide, "Définition des 5 endpoints REST de l'API", 0.5, 1.3, 12, 0.4, size=14, color=GRAY)

rect(slide, 0.5, 1.85, 6.5, 3.2, RGBColor(0x1E, 0x1E, 0x2E))
text_box(slide,
    "Route::middleware('auth')->group(\n"
    "  function () {\n\n"
    "    // NOTES\n"
    "    Route::get('/notes',\n"
    "      [NoteController::class, 'index']);\n"
    "    Route::post('/notes',\n"
    "      [NoteController::class, 'store']);\n"
    "    Route::delete('/notes/{id}',\n"
    "      [NoteController::class, 'destroy']);\n\n"
    "    // TAGS\n"
    "    Route::get('/tags',\n"
    "      [TagController::class, 'index']);\n"
    "    Route::post('/tags',\n"
    "      [TagController::class, 'store']);\n"
    "  }\n"
    ");",
    0.65, 1.95, 6.2, 3.0, size=11, color=ORANGE)

text_box(slide, "Tableau des endpoints :", 7.3, 1.85, 5.8, 0.4, size=14, bold=True, color=WHITE)
endpoints = [
    ("GET",    "/api/notes",       "Liste les notes"),
    ("POST",   "/api/notes",       "Crée une note"),
    ("DELETE", "/api/notes/{id}",  "Supprime une note"),
    ("GET",    "/api/tags",        "Liste les tags"),
    ("POST",   "/api/tags",        "Crée un tag"),
]
method_colors = {"GET": GREEN, "POST": ACCENT, "DELETE": RED}
for i, (method, path, desc) in enumerate(endpoints):
    y = 2.35 + i * 0.72
    col = method_colors.get(method, YELLOW)
    rect(slide, 7.3, y, 1.2, 0.55, RGBColor(0x31, 0x32, 0x44))
    text_box(slide, method, 7.35, y+0.08, 1.1, 0.4, size=12, bold=True, color=col)
    text_box(slide, path,   8.65, y+0.08, 2.5, 0.4, size=11, color=WHITE)
    text_box(slide, desc,   11.2, y+0.08, 2.0, 0.4, size=11, color=GRAY)

text_box(slide, "⚠  Tous protégés par middleware 'auth' — accès refusé sans session active.",
         0.5, 5.2, 12, 0.5, size=13, color=YELLOW)
text_box(slide, "✅  bootstrap/app.php modifié pour enregistrer api.php sur le préfixe /api",
         0.5, 5.8, 12, 0.5, size=13, color=GREEN)


# ── SLIDE 13 : Avant / Après — comparaison ────────────────────────────────────
slide = add_slide()
text_box(slide, "Comparaison Avant / Après refactorisation", 0.5, 0.3, 12, 0.8, size=26, bold=True, color=ACCENT)
rect(slide, 0.5, 1.1, 12.3, 0.04, ACCENT)

# Avant
text_box(slide, "AVANT", 0.5, 1.3, 6, 0.45, size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
rect(slide, 0.5, 1.8, 6.0, 4.8, RGBColor(0x31, 0x32, 0x44))
avant_items = [
    "❌  Notes.php : validation + BDD + UI + flash",
    "❌  TagForm.php : validation + création + events",
    "❌  Aucune route API (/api/*)",
    "❌  HTML généré côté serveur (Livewire)",
    "❌  Impossible de tester sans UI",
    "❌  Tag::all() dupliqué dans 2 méthodes",
    "❌  Front non séparable du back",
]
for i, item in enumerate(avant_items):
    text_box(slide, item, 0.65, 1.95 + i*0.62, 5.7, 0.55, size=13, color=RED if "❌" in item else WHITE)

# Après
text_box(slide, "APRÈS", 6.8, 1.3, 6, 0.45, size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
rect(slide, 6.8, 1.8, 6.0, 4.8, RGBColor(0x31, 0x32, 0x44))
apres_items = [
    "✅  NoteService : logique métier isolée",
    "✅  TagService : logique métier isolée",
    "✅  5 endpoints REST dans routes/api.php",
    "✅  Réponses JSON uniformes { status, data }",
    "✅  NoteController : reçoit HTTP, délègue, répond",
    "✅  Injection de dépendances (principe D)",
    "✅  Prêt pour front React découplé",
]
for i, item in enumerate(apres_items):
    text_box(slide, item, 6.95, 1.95 + i*0.62, 5.7, 0.55, size=13, color=GREEN if "✅" in item else WHITE)

# Séparateur
rect(slide, 6.6, 1.3, 0.06, 5.3, GRAY)


# ── SLIDE 14 : Pourquoi React ? ───────────────────────────────────────────────
slide = add_slide()
text_box(slide, "Pourquoi migrer vers React ?", 0.5, 0.3, 12, 0.8, size=28, bold=True, color=ACCENT)
rect(slide, 0.5, 1.1, 12.3, 0.04, ACCENT)

# Livewire vs React
text_box(slide, "LIVEWIRE (avant)", 0.5, 1.4, 6.0, 0.45, size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
text_box(slide, "REACT (après)", 6.8, 1.4, 6.0, 0.45, size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
rect(slide, 6.6, 1.3, 0.06, 5.5, GRAY)

livewire_points = [
    "HTML généré côté SERVEUR PHP",
    "Chaque interaction = aller-retour serveur",
    "wire:model / wire:submit → couplé PHP",
    "Impossible d'avoir app mobile",
    "Impossible d'avoir 2ème front",
    "Tests UI = tests PHP complets",
    "1 déploiement = tout en PHP",
]
react_points = [
    "HTML géré côté NAVIGATEUR",
    "Interactions instantanées (pas de rechargement)",
    "Composants JSX réutilisables",
    "Le même back peut servir web + mobile",
    "N'importe quel front peut consommer l'API",
    "Tests UI séparés des tests PHP",
    "Back et front déployables indépendamment",
]
for i, (lw, rc) in enumerate(zip(livewire_points, react_points)):
    y = 1.95 + i * 0.6
    text_box(slide, "✗  " + lw, 0.6, y, 5.8, 0.5, size=12, color=RED)
    text_box(slide, "✓  " + rc, 6.85, y, 5.8, 0.5, size=12, color=GREEN)


# ── SLIDE 15 : Zustand — pourquoi ce choix ────────────────────────────────────
slide = add_slide()
text_box(slide, "State Management — Pourquoi Zustand ?", 0.5, 0.3, 12, 0.8, size=26, bold=True, color=ACCENT)
rect(slide, 0.5, 1.1, 12.3, 0.04, ACCENT)

text_box(slide, "Le problème sans state management :", 0.5, 1.3, 12, 0.4, size=14, bold=True, color=YELLOW)
add_bullet(slide, [
    "NoteForm et NoteList ont tous deux besoin des mêmes données (notes, tags)",
    "Sans store partagé → il faudrait remonter l'état jusqu'au parent commun (prop drilling)",
    "Exemple : TagForm crée un tag → NoteForm doit le voir → sans store, il faudrait tout faire remonter",
], 0.5, 1.75, 12, 1.2, size=13, color=GRAY)

text_box(slide, "Comparaison des solutions :", 0.5, 3.1, 12, 0.4, size=14, bold=True, color=WHITE)
options = [
    ("Redux",   "Très puissant\nBoilerplate important\nParfait pour grandes apps",      RED,    "❌ Trop complexe pour ce projet"),
    ("Zustand", "Simple, léger\nPas de boilerplate\nAPI intuitive (create + hooks)",    GREEN,  "✅ CHOIX RETENU"),
    ("MobX",    "Basé sur MVVM\nObservables automatiques\nCourbe d'apprentissage +",   YELLOW, "⚠ Bien mais + complexe"),
    ("Context", "Natif React\nPas de lib externe\nRe-renders excessifs si mal utilisé",ACCENT, "⚠ Suffisant mais limité"),
]
for i, (name, desc, col, verdict) in enumerate(options):
    bx = 0.4 + i * 3.2
    rect(slide, bx, 3.6, 3.0, 2.5, RGBColor(0x31, 0x32, 0x44))
    rect(slide, bx, 3.6, 3.0, 0.35, col)
    text_box(slide, name, bx+0.1, 3.62, 2.8, 0.3, size=14, bold=True, color=DARK_BG)
    text_box(slide, desc, bx+0.1, 4.05, 2.8, 1.3, size=11, color=WHITE)
    text_box(slide, verdict, bx+0.1, 5.5, 2.8, 0.45, size=11, bold=True, color=col)


# ── SLIDE 16 : useAppStore.js — expliqué ─────────────────────────────────────
slide = add_slide()
text_box(slide, "store/useAppStore.js — Le cœur de l'état React", 0.5, 0.3, 12, 0.8, size=24, bold=True, color=ACCENT)
rect(slide, 0.5, 1.1, 12.3, 0.04, ACCENT)

rect(slide, 0.5, 1.3, 7.5, 5.3, RGBColor(0x1E, 0x1E, 0x2E))
text_box(slide,
    "const useAppStore = create((set) => ({\n\n"
    "  // État global partagé\n"
    "  notes:   [],\n"
    "  tags:    [],\n"
    "  loading: false,\n"
    "  error:   null,\n\n"
    "  // Action : charge les notes depuis l'API\n"
    "  fetchNotes: async () => {\n"
    "    set({ loading: true });\n"
    "    const notes = await getNotes(); // ← api.js\n"
    "    set({ notes, loading: false });\n"
    "  },\n\n"
    "  // Action : crée une note + met à jour le store\n"
    "  addNote: async (text, tag_id) => {\n"
    "    const note = await createNote(text, tag_id);\n"
    "    set((state) => ({\n"
    "      notes: [note, ...state.notes]\n"
    "    }));\n"
    "  },\n"
    "}));",
    0.65, 1.4, 7.2, 5.1, size=10.5, color=ORANGE)

annotations = [
    (1.6,  "État global : partagé entre NoteForm, NoteList, TagForm"),
    (3.05, "Délègue à api.js — le store ne connaît pas fetch/axios"),
    (3.65, "set() met à jour l'état → tous les composants qui lisent\nnotes se re-rendent automatiquement"),
    (4.65, "Optimistic update : on ajoute au store sans re-fetcher"),
    (5.35, "Pattern Flux : action → store → UI"),
]
for ay, txt in annotations:
    rect(slide, 8.2, ay, 5.0, 0.65, RGBColor(0x31, 0x32, 0x44))
    rect(slide, 8.2, ay, 0.06, 0.65, ACCENT)
    text_box(slide, txt, 8.35, ay+0.08, 4.7, 0.55, size=11, color=WHITE)


# ── SLIDE 17 : Flux de données complet ───────────────────────────────────────
slide = add_slide()
text_box(slide, "Flux de données complet — Exemple : créer une note", 0.5, 0.3, 12, 0.8, size=24, bold=True, color=YELLOW)
rect(slide, 0.5, 1.1, 12.3, 0.04, YELLOW)

steps = [
    (ACCENT,  "1. Utilisateur remplit le formulaire",
               "NoteForm.jsx — state local : setText(), setTagId()"),
    (YELLOW,  "2. Clic sur 'Ajouter la note'",
               "handleSubmit() → appelle addNote(text, tag_id) du store Zustand"),
    (ORANGE,  "3. Store appelle l'API",
               "useAppStore.addNote() → createNote() dans api.js → POST /api/notes"),
    (GREEN,   "4. Laravel reçoit la requête",
               "routes/api.php → NoteController.store() → NoteService.create() → BDD SQLite"),
    (GREEN,   "5. Réponse JSON 201",
               '{ "status": "success", "data": { id, text, tag, ... } }'),
    (ACCENT,  "6. Store met à jour l'état",
               "set({ notes: [newNote, ...state.notes] }) → NoteList se re-rend"),
    (WHITE,   "7. UI mise à jour sans rechargement",
               "NoteList affiche la nouvelle note en tête de liste"),
]
for i, (col, title, detail) in enumerate(steps):
    y = 1.4 + i * 0.72
    rect(slide, 0.4, y, 0.5, 0.58, col)
    text_box(slide, str(i+1), 0.45, y+0.08, 0.4, 0.42, size=16, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    text_box(slide, title,  1.05, y+0.05, 5.5, 0.48, size=13, bold=True, color=col)
    text_box(slide, detail, 6.7,  y+0.05, 6.4, 0.48, size=12, color=GRAY)
    if i < len(steps) - 1:
        text_box(slide, "↓", 0.55, y+0.6, 0.4, 0.2, size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 18 : Structure finale des fichiers ──────────────────────────────────
slide = add_slide()
text_box(slide, "Structure finale du projet", 0.5, 0.3, 12, 0.8, size=28, bold=True, color=ACCENT)
rect(slide, 0.5, 1.1, 12.3, 0.04, ACCENT)

text_box(slide, "BACK-END PHP/Laravel", 0.5, 1.3, 6.0, 0.4, size=14, bold=True, color=YELLOW)
back_files = [
    ("app/Services/NoteService.php",                  "✅ NOUVEAU — logique métier notes",       GREEN),
    ("app/Services/TagService.php",                   "✅ NOUVEAU — logique métier tags",        GREEN),
    ("app/Http/Controllers/Api/NoteController.php",   "✅ NOUVEAU — contrôleur API notes",       GREEN),
    ("app/Http/Controllers/Api/TagController.php",    "✅ NOUVEAU — contrôleur API tags",        GREEN),
    ("routes/api.php",                                "✅ NOUVEAU — 5 endpoints REST",           GREEN),
    ("app/Models/Note.php + Tag.php + User.php",      "✅ CONSERVÉ — modèles Eloquent",          ACCENT),
    ("bootstrap/app.php",                             "⚠  MODIFIÉ — enregistre api.php",        YELLOW),
    ("app/Livewire/Notes.php",                        "🗑  OBSOLÈTE — remplacé par NoteService", RED),
    ("app/Livewire/TagForm.php",                      "🗑  OBSOLÈTE — remplacé par TagService",  RED),
]
for i, (path, label, col) in enumerate(back_files):
    y = 1.75 + i * 0.53
    text_box(slide, path,  0.5, y, 5.5, 0.45, size=10, color=WHITE)
    text_box(slide, label, 6.1, y, 4.5, 0.45, size=10, color=col)

text_box(slide, "FRONT-END React", 6.8, 1.3, 6.0, 0.4, size=14, bold=True, color=ACCENT)
front_files = [
    ("resources/js/app.jsx",                "✅ NOUVEAU — point d'entrée React"),
    ("resources/js/components/App.jsx",     "✅ NOUVEAU — composant racine"),
    ("resources/js/components/NoteForm.jsx","✅ NOUVEAU — formulaire note"),
    ("resources/js/components/NoteList.jsx","✅ NOUVEAU — liste des notes"),
    ("resources/js/components/TagForm.jsx", "✅ NOUVEAU — formulaire tag"),
    ("resources/js/store/useAppStore.js",   "✅ NOUVEAU — store Zustand"),
    ("resources/js/services/api.js",        "✅ NOUVEAU — appels API REST"),
    ("vite.config.js",                      "⚠  MODIFIÉ — plugin React JSX"),
    ("views/dashboard.blade.php",           "⚠  MODIFIÉ — monte React"),
]
for i, (path, label) in enumerate(front_files):
    y = 1.75 + i * 0.53
    col = GREEN if "NOUVEAU" in label else YELLOW
    text_box(slide, path,  6.8, y, 4.0, 0.45, size=10, color=WHITE)
    text_box(slide, label, 11.0, y, 2.2, 0.45, size=9, color=col)

rect(slide, 6.6, 1.3, 0.06, 5.3, GRAY)


# ── SLIDE 19 : Les 5 principes SOLID — vue d'ensemble ───────────────────────
slide = add_slide()
text_box(slide, "Les 5 principes SOLID — rappel complet", 0.5, 0.3, 12, 0.7, size=28, bold=True, color=ACCENT)
rect(slide, 0.5, 1.05, 12.3, 0.04, ACCENT)
text_box(slide, "Applicables à tout langage OOP — PHP, Java, C#, Python…", 0.5, 1.15, 12, 0.4, size=14, color=GRAY)

solids = [
    ("S", "Single Responsibility",  "Une classe = une seule raison de changer",
     "NoteService ne fait que la logique notes\nNoteController ne fait que gérer HTTP\nNote (Model) ne fait qu'accéder à la BDD",
     GREEN),
    ("O", "Open / Closed",          "Ouvert à l'extension, fermé à la modification",
     "Ajouter un TagService n'oblige pas à modifier NoteService\nAjouter un endpoint n'oblige pas à modifier les autres contrôleurs",
     ACCENT),
    ("L", "Liskov Substitution",    "Une sous-classe doit pouvoir remplacer sa classe parente",
     "Non directement appliqué ici\nMais respecté : NoteController extends Controller sans casser les comportements",
     YELLOW),
    ("I", "Interface Segregation",  "Ne pas forcer une classe à implémenter ce dont elle n'a pas besoin",
     "NoteController n'expose que les routes notes\nTagController n'expose que les routes tags\nPas de méthode inutile dans les services",
     ORANGE),
    ("D", "Dependency Inversion",   "Dépendre des abstractions, pas des implémentations concrètes",
     "NoteController reçoit NoteService via constructeur\nLaravel injecte automatiquement — pas de new NoteService()\napi.js est l'abstraction pour les composants React",
     RED),
]
for i, (letter, name, definition, example, col) in enumerate(solids):
    bx = 0.35 + i * 2.52
    rect(slide, bx, 1.7, 2.42, 5.4, RGBColor(0x31, 0x32, 0x44))
    rect(slide, bx, 1.7, 2.42, 0.5, col)
    text_box(slide, letter, bx+0.08, 1.72, 0.4, 0.45, size=22, bold=True, color=DARK_BG)
    text_box(slide, name,   bx+0.5,  1.75, 1.85, 0.45, size=11, bold=True, color=DARK_BG)
    text_box(slide, definition, bx+0.1, 2.3, 2.25, 0.7, size=10, color=WHITE)
    rect(slide, bx+0.1, 3.05, 2.25, 0.04, col)
    text_box(slide, "Dans Renote :", bx+0.1, 3.15, 2.25, 0.3, size=9, bold=True, color=col)
    text_box(slide, example, bx+0.1, 3.5, 2.25, 2.5, size=9, color=GRAY)


# ── SLIDE 20 : MVC — pattern expliqué avec Renote ────────────────────────────
slide = add_slide()
text_box(slide, "Le pattern MVC — comment il s'applique à Renote", 0.5, 0.3, 12, 0.7, size=26, bold=True, color=YELLOW)
rect(slide, 0.5, 1.05, 12.3, 0.04, YELLOW)

# Schéma MVC
layers_mvc = [
    ("M  MODEL",      "app/Models/Note.php\napp/Models/Tag.php\napp/Models/User.php",
     "Représente les données et les règles métier.\nParle directement à la base de données via Eloquent ORM.\nDéfinit les relations (belongsTo, hasMany).",
     ACCENT,   0.4,  1.4),
    ("V  VIEW",       "resources/js/components/\nApp.jsx · NoteForm.jsx\nNoteList.jsx · TagForm.jsx",
     "Affiche les données à l'utilisateur.\nNe contient aucune logique métier.\nObserve le store Zustand et se re-rend automatiquement.",
     GREEN,    0.4,  3.5),
    ("C  CONTROLLER", "app/Http/Controllers/Api/\nNoteController.php\nTagController.php",
     "Fait le lien entre Model et View.\nReçoit la requête HTTP, valide les données,\ndélègue au Service, renvoie la réponse JSON.",
     ORANGE,   0.4,  5.5),
]
for m_letter, files, desc, col, bx, by in layers_mvc:
    rect(slide, bx, by, 3.5, 1.7, RGBColor(0x31, 0x32, 0x44))
    rect(slide, bx, by, 0.12, 1.7, col)
    text_box(slide, m_letter, bx+0.25, by+0.1, 3.1, 0.4, size=13, bold=True, color=col)
    text_box(slide, files,    bx+0.25, by+0.55, 3.1, 0.7, size=10, color=GRAY)

# Flèches
text_box(slide, "↕  queries / data", 0.7, 3.2, 3.0, 0.3, size=10, color=GRAY)
text_box(slide, "↕  JSON responses", 0.7, 5.2, 3.0, 0.3, size=10, color=GRAY)

# Explication côté droit
rect(slide, 4.4, 1.4, 8.8, 5.8, RGBColor(0x31, 0x32, 0x44))
text_box(slide, "Pourquoi MVC ?", 4.6, 1.5, 8.4, 0.45, size=16, bold=True, color=YELLOW)

mvc_points = [
    ("Séparation des responsabilités",
     "Chaque couche a un rôle clair. Changer le design (View)\nn'oblige pas à modifier la logique (Controller/Model)."),
    ("Maintenabilité",
     "Un bug dans l'affichage → on regarde les composants React.\nUn bug de données → on regarde les Models/Services."),
    ("Testabilité",
     "On peut tester NoteService sans lancer Laravel.\nOn peut tester NoteList sans avoir de back-end."),
    ("Extension facile",
     "Ajouter un endpoint = créer une méthode dans le Controller.\nAjouter une relation = modifier uniquement le Model."),
    ("Note sur les Services",
     "Dans Renote on ajoute une couche Service entre Controller et Model.\nC'est du MVC+ (ou MVC avec Service Layer) — meilleure pratique SOLID."),
]
for i, (title, desc) in enumerate(mvc_points):
    y = 2.05 + i * 0.98
    rect(slide, 4.5, y, 0.06, 0.8, YELLOW)
    text_box(slide, title, 4.7, y,      8.2, 0.35, size=12, bold=True, color=WHITE)
    text_box(slide, desc,  4.7, y+0.38, 8.2, 0.55, size=11, color=GRAY)


# ── SLIDE 21 : Observer pattern → origine du MVC ─────────────────────────────
slide = add_slide()
text_box(slide, "Le pattern Observer — 'papa' du MVC", 0.5, 0.3, 12, 0.7, size=26, bold=True, color=ORANGE)
rect(slide, 0.5, 1.05, 12.3, 0.04, ORANGE)

text_box(slide,
    "Le pattern Observer définit une relation 1-à-N entre objets :\nquand l'objet observé (Subject) change d'état, tous ses observateurs (Observers) sont notifiés automatiquement.",
    0.5, 1.2, 12, 0.7, size=14, color=GRAY)

# Schéma Observer pur
rect(slide, 0.5, 2.0, 3.5, 1.2, RGBColor(0x31, 0x32, 0x44))
text_box(slide, "SUBJECT (Observable)", 0.65, 2.05, 3.2, 0.4, size=12, bold=True, color=ORANGE)
text_box(slide, "setState()\nnotifyObservers()", 0.65, 2.5, 3.2, 0.6, size=11, color=GRAY)

for i, obs in enumerate(["Observer A", "Observer B", "Observer C"]):
    bx = 5.0 + i * 2.7
    rect(slide, bx, 2.0, 2.3, 1.2, RGBColor(0x31, 0x32, 0x44))
    text_box(slide, obs, bx+0.15, 2.15, 2.0, 0.4, size=12, bold=True, color=ACCENT)
    text_box(slide, "update()", bx+0.15, 2.6, 2.0, 0.4, size=11, color=GRAY)
    text_box(slide, "─────────→", 4.1 + i*2.7, 2.45, 1.0, 0.35, size=13, color=ORANGE)

text_box(slide, "Observer pur", 0.5, 3.3, 12, 0.35, size=11, color=GRAY, align=PP_ALIGN.CENTER)
rect(slide, 0.5, 3.65, 12.3, 0.04, ORANGE)

# Lien avec MVC
text_box(slide, "Dans MVC, l'Observer devient :", 0.5, 3.8, 12, 0.4, size=14, bold=True, color=WHITE)

mvc_obs = [
    ("MODEL",      "= Subject (Observable)\nContient les données.\nNotifie quand les données changent.",     ACCENT),
    ("VIEW",       "= Observer\nS'abonne au Model.\nSe met à jour automatiquement quand le Model change.",   GREEN),
    ("CONTROLLER", "= Intermédiaire\nReçoit les actions utilisateur.\nMet à jour le Model.",                 ORANGE),
]
for i, (role, desc, col) in enumerate(mvc_obs):
    bx = 0.5 + i * 4.2
    rect(slide, bx, 4.3, 3.9, 2.5, RGBColor(0x31, 0x32, 0x44))
    rect(slide, bx, 4.3, 3.9, 0.4, col)
    text_box(slide, role, bx+0.15, 4.33, 3.6, 0.35, size=13, bold=True, color=DARK_BG)
    text_box(slide, desc, bx+0.15, 4.8,  3.6, 1.8,  size=11, color=WHITE)

# Dans Renote
text_box(slide, "Dans Renote : Zustand joue le rôle du Subject Observer — les composants React sont les Observers",
         0.5, 6.95, 12, 0.45, size=13, color=YELLOW, align=PP_ALIGN.CENTER)


# ── SLIDE 22 : REST — conventions complètes ──────────────────────────────────
slide = add_slide()
text_box(slide, "Conventions REST — comment on les applique dans Renote", 0.5, 0.3, 12, 0.7, size=24, bold=True, color=ACCENT)
rect(slide, 0.5, 1.05, 12.3, 0.04, ACCENT)

# Verbes HTTP
text_box(slide, "Les verbes HTTP et leur sémantique :", 0.5, 1.2, 12, 0.4, size=14, bold=True, color=WHITE)
verbs = [
    ("GET",    "Lecture seule — ne modifie rien",      "/api/notes → liste toutes les notes",           GREEN),
    ("POST",   "Création d'une ressource",             "/api/notes → crée une note (body JSON requis)", ACCENT),
    ("PUT",    "Remplacement complet d'une ressource", "/api/notes/1 → remplace la note 1 (non utilisé ici)", YELLOW),
    ("PATCH",  "Modification partielle",               "/api/notes/1 → modifie le texte seulement (non utilisé ici)", YELLOW),
    ("DELETE", "Suppression d'une ressource",          "/api/notes/1 → supprime la note 1",             RED),
]
for i, (verb, meaning, example, col) in enumerate(verbs):
    y = 1.7 + i * 0.58
    rect(slide, 0.5, y, 1.4, 0.48, RGBColor(0x31, 0x32, 0x44))
    text_box(slide, verb, 0.55, y+0.07, 1.3, 0.38, size=13, bold=True, color=col)
    text_box(slide, meaning,  2.1, y+0.07, 4.5, 0.38, size=12, color=WHITE)
    text_box(slide, example,  6.8, y+0.07, 6.0, 0.38, size=11, color=GRAY)

rect(slide, 0.5, 4.65, 12.3, 0.04, ACCENT)

# Codes HTTP
text_box(slide, "Codes HTTP de réponse utilisés dans l'API :", 0.5, 4.8, 12, 0.4, size=14, bold=True, color=WHITE)
codes = [
    ("200 OK",              "Requête réussie (GET, DELETE)",              GREEN),
    ("201 Created",         "Ressource créée (POST réussi)",              GREEN),
    ("400 Bad Request",     "Données invalides (validation échouée)",     YELLOW),
    ("401 Unauthorized",    "Non authentifié (pas de session)",           RED),
    ("404 Not Found",       "Ressource inexistante ou accès refusé",      RED),
    ("422 Unprocessable",   "Validation Laravel — champ manquant/invalide",ORANGE),
]
for i, (code, desc, col) in enumerate(codes):
    bx = 0.5 + (i % 3) * 4.2
    by = 5.3 + (i // 3) * 0.7
    rect(slide, bx, by, 1.5, 0.55, col)
    text_box(slide, code, bx+0.1, by+0.08, 1.3, 0.42, size=10, bold=True, color=DARK_BG)
    text_box(slide, desc, bx+1.65, by+0.08, 2.4, 0.42, size=10, color=WHITE)


# ── SLIDE 23 : Injection de dépendances — deep dive ──────────────────────────
slide = add_slide()
text_box(slide, "Injection de dépendances — principe D en profondeur", 0.5, 0.3, 12, 0.7, size=24, bold=True, color=RED)
rect(slide, 0.5, 1.05, 12.3, 0.04, RED)

# Sans injection
text_box(slide, "❌ SANS injection (couplage fort)", 0.5, 1.2, 6.0, 0.4, size=14, bold=True, color=RED)
rect(slide, 0.5, 1.65, 5.8, 2.5, RGBColor(0x1E, 0x1E, 0x2E))
text_box(slide,
    "class NoteController\n"
    "{\n"
    "  public function index()\n"
    "  {\n"
    "    // ❌ Crée lui-même sa dépendance\n"
    "    $service = new NoteService();\n"
    "    $notes = $service->getAllForUser();\n"
    "    ...\n"
    "  }\n"
    "}",
    0.65, 1.75, 5.5, 2.3, size=11, color=RED)

# Avec injection
text_box(slide, "✅ AVEC injection (couplage faible)", 7.0, 1.2, 6.0, 0.4, size=14, bold=True, color=GREEN)
rect(slide, 7.0, 1.65, 5.8, 2.5, RGBColor(0x1E, 0x1E, 0x2E))
text_box(slide,
    "class NoteController\n"
    "{\n"
    "  // ✅ Laravel injecte NoteService\n"
    "  // automatiquement\n"
    "  public function __construct(\n"
    "    private NoteService $noteService\n"
    "  ) {}\n\n"
    "  public function index()\n"
    "  {\n"
    "    $notes = $this->noteService\n"
    "               ->getAllForUser();\n"
    "  }\n"
    "}",
    7.15, 1.75, 5.5, 2.3, size=11, color=GREEN)

# Séparateur
rect(slide, 6.6, 1.2, 0.06, 2.95, GRAY)

# Pourquoi c'est important
rect(slide, 0.5, 4.3, 12.3, 0.04, RED)
text_box(slide, "Pourquoi l'injection de dépendances est cruciale :", 0.5, 4.45, 12, 0.4, size=14, bold=True, color=WHITE)

benefits = [
    ("Testabilité",
     "Sans injection : impossible de tester NoteController sans une vraie BDD.\n"
     "Avec injection : on peut passer un NoteService 'fake' dans les tests unitaires."),
    ("Flexibilité",
     "Demain on veut changer NoteService par NoteServiceV2 ?\n"
     "On change une seule ligne (le binding dans AppServiceProvider), pas chaque contrôleur."),
    ("Lisibilité",
     "Le constructeur liste explicitement toutes les dépendances d'une classe.\n"
     "On sait immédiatement de quoi NoteController a besoin pour fonctionner."),
]
for i, (title, desc) in enumerate(benefits):
    bx = 0.5 + i * 4.2
    rect(slide, bx, 4.95, 3.9, 2.2, RGBColor(0x31, 0x32, 0x44))
    rect(slide, bx, 4.95, 3.9, 0.35, RED)
    text_box(slide, title, bx+0.15, 4.97, 3.6, 0.3, size=12, bold=True, color=DARK_BG)
    text_box(slide, desc,  bx+0.15, 5.38, 3.6, 1.65, size=10, color=WHITE)


# ── SLIDE 24 : TagService + TagController — code annoté ──────────────────────
slide = add_slide()
text_box(slide, "TagService.php + TagController.php — vue complète", 0.5, 0.3, 12, 0.7, size=24, bold=True, color=GREEN)
rect(slide, 0.5, 1.05, 12.3, 0.04, GREEN)

# TagService
text_box(slide, "app/Services/TagService.php", 0.5, 1.2, 6.0, 0.35, size=12, color=GRAY)
rect(slide, 0.5, 1.6, 5.8, 3.2, RGBColor(0x1E, 0x1E, 0x2E))
text_box(slide,
    "class TagService\n"
    "{\n"
    "  // Retourne tous les tags (globaux)\n"
    "  // Pas lié à un user — tags partagés\n"
    "  public function getAll()\n"
    "  {\n"
    "    return Tag::all();\n"
    "  }\n\n"
    "  // Crée un tag\n"
    "  // Validation faite en amont\n"
    "  // dans le Controller\n"
    "  public function create(array $data)\n"
    "  {\n"
    "    return Tag::create(['name' => $data['name']]);\n"
    "  }\n"
    "}",
    0.65, 1.7, 5.5, 3.0, size=11, color=GREEN)

# TagController
text_box(slide, "app/Http/Controllers/Api/TagController.php", 6.5, 1.2, 6.5, 0.35, size=12, color=GRAY)
rect(slide, 6.5, 1.6, 6.3, 3.2, RGBColor(0x1E, 0x1E, 0x2E))
text_box(slide,
    "class TagController extends Controller\n"
    "{\n"
    "  // Injection automatique Laravel\n"
    "  public function __construct(\n"
    "    private TagService $tagService\n"
    "  ) {}\n\n"
    "  // GET /api/tags\n"
    "  public function index(): JsonResponse\n"
    "  {\n"
    "    return response()->json([\n"
    "      'status' => 'success',\n"
    "      'data'   => $this->tagService->getAll(),\n"
    "    ]);\n"
    "  }\n\n"
    "  // POST /api/tags\n"
    "  public function store(Request $req)\n"
    "  {\n"
    "    $v = $req->validate([\n"
    "      'name' => 'required|max:50|unique:tags,name'\n"
    "    ]);\n"
    "    return response()->json([\n"
    "      'data' => $this->tagService->create($v)\n"
    "    ], 201);\n"
    "  }\n"
    "}",
    6.65, 1.7, 6.0, 3.0, size=10, color=ORANGE)

# Points clés
rect(slide, 0.5, 4.95, 12.3, 0.04, GREEN)
text_box(slide, "Points clés à retenir :", 0.5, 5.1, 12, 0.35, size=13, bold=True, color=WHITE)
pts = [
    "TagService.getAll() ne filtre pas par user — les tags sont globaux (contrairement aux notes)",
    "La validation 'unique:tags,name' est dans le Controller, pas dans le Service (frontière système)",
    "Même pattern que NoteController : injection → validation → délègue → répond JSON",
    "Les deux contrôleurs héritent de Controller (base Laravel) — accès aux helpers response(), etc.",
]
for i, pt in enumerate(pts):
    bx = 0.5 + (i % 2) * 6.2
    by = 5.55 + (i // 2) * 0.6
    text_box(slide, "▸  " + pt, bx, by, 6.0, 0.5, size=11, color=GRAY)


# ── SLIDE 25 : api.js — service layer React ───────────────────────────────────
slide = add_slide()
text_box(slide, "services/api.js — La couche d'accès API côté React", 0.5, 0.3, 12, 0.7, size=24, bold=True, color=ACCENT)
rect(slide, 0.5, 1.05, 12.3, 0.04, ACCENT)

rect(slide, 0.5, 1.2, 6.2, 5.5, RGBColor(0x1E, 0x1E, 0x2E))
text_box(slide,
    "// Instance axios partagée\n"
    "const api = axios.create({\n"
    "  baseURL: '/api',\n"
    "  withCredentials: true,  // cookie session\n"
    "  headers: {\n"
    "    'X-CSRF-TOKEN': document\n"
    "      .querySelector('meta[name=csrf-token]')\n"
    "      ?.content,\n"
    "  },\n"
    "});\n\n"
    "// GET /api/notes\n"
    "export const getNotes = async () => {\n"
    "  const res = await api.get('/notes');\n"
    "  return res.data.data;\n"
    "};\n\n"
    "// POST /api/notes\n"
    "export const createNote = async (text, tag_id) => {\n"
    "  const res = await api.post('/notes', { text, tag_id });\n"
    "  return res.data.data;\n"
    "};\n\n"
    "// DELETE /api/notes/{id}\n"
    "export const deleteNote = async (id) => {\n"
    "  await api.delete(`/notes/${id}`);\n"
    "};",
    0.65, 1.3, 5.9, 5.3, size=10, color=ORANGE)

annotations_api = [
    (1.35, "Principe D : les composants dépendent\nde ces fonctions, pas de axios directement"),
    (2.05, "withCredentials: true → envoie le cookie\nde session Laravel avec chaque requête"),
    (2.7,  "CSRF token lu dans le HTML Blade → obligatoire\npour POST/DELETE sinon Laravel rejette (419)"),
    (3.55, "Chaque fonction = un seul endpoint\nPrincipe S : une responsabilité"),
    (4.45, "res.data.data → on extrait data du JSON\n{ status, data: [...] }"),
    (5.55, "Même pattern pour tags :\ngetTags(), createTag(name)"),
]
for ay, txt in annotations_api:
    rect(slide, 6.9, ay, 6.0, 0.72, RGBColor(0x31, 0x32, 0x44))
    rect(slide, 6.9, ay, 0.07, 0.72, ACCENT)
    text_box(slide, txt, 7.05, ay+0.08, 5.75, 0.6, size=11, color=WHITE)


# ── SLIDE 26 : NoteForm + NoteList — composants React annotés ────────────────
slide = add_slide()
text_box(slide, "Composants React — NoteForm.jsx + NoteList.jsx", 0.5, 0.3, 12, 0.7, size=24, bold=True, color=GREEN)
rect(slide, 0.5, 1.05, 12.3, 0.04, GREEN)

# NoteForm
text_box(slide, "NoteForm.jsx", 0.5, 1.2, 6.0, 0.35, size=12, bold=True, color=GREEN)
rect(slide, 0.5, 1.6, 5.8, 4.0, RGBColor(0x1E, 0x1E, 0x2E))
text_box(slide,
    "export default function NoteForm() {\n"
    "  // État LOCAL du form (non partagé)\n"
    "  const [text, setText] = useState('');\n"
    "  const [tagId, setTagId] = useState('');\n\n"
    "  // Sélecteurs ciblés dans le store\n"
    "  // → re-render seulement si ces valeurs changent\n"
    "  const tags    = useAppStore(s => s.tags);\n"
    "  const addNote = useAppStore(s => s.addNote);\n"
    "  const loading = useAppStore(s => s.loading);\n\n"
    "  const handleSubmit = async (e) => {\n"
    "    e.preventDefault();\n"
    "    // Délègue au store, pas de fetch ici\n"
    "    await addNote(text, parseInt(tagId));\n"
    "    setText(''); setTagId(''); // reset\n"
    "  };\n\n"
    "  return <form onSubmit={handleSubmit}>...</form>;\n"
    "}",
    0.65, 1.7, 5.5, 3.8, size=10, color=ORANGE)

# NoteList
text_box(slide, "NoteList.jsx", 6.5, 1.2, 6.5, 0.35, size=12, bold=True, color=ACCENT)
rect(slide, 6.5, 1.6, 6.3, 4.0, RGBColor(0x1E, 0x1E, 0x2E))
text_box(slide,
    "export default function NoteList() {\n"
    "  const notes      = useAppStore(s => s.notes);\n"
    "  const loading    = useAppStore(s => s.loading);\n"
    "  const fetchNotes = useAppStore(s => s.fetchNotes);\n"
    "  const removeNote = useAppStore(s => s.removeNote);\n\n"
    "  // useEffect = componentDidMount\n"
    "  // [] = s'exécute 1 seule fois au montage\n"
    "  useEffect(() => {\n"
    "    fetchNotes(); // → GET /api/notes\n"
    "  }, []);\n\n"
    "  return (\n"
    "    <div>\n"
    "      {notes.map(note => (\n"
    "        <div key={note.id}>\n"
    "          <p>{note.text}</p>\n"
    "          <span>{note.tag?.name}</span>\n"
    "          <button onClick={() => removeNote(note.id)}>\n"
    "            Supprimer\n"
    "          </button>\n"
    "        </div>\n"
    "      ))}\n"
    "    </div>\n"
    "  );\n"
    "}",
    6.65, 1.7, 6.0, 3.8, size=10, color=ORANGE)

rect(slide, 0.5, 5.75, 12.3, 0.04, GREEN)
pts2 = [
    "État local vs store : text/tagId sont locaux (form) — notes/tags sont globaux (partagés entre composants)",
    "note.tag?.name → opérateur ?. (optional chaining) : si tag est null, n'explose pas",
    "useEffect(fn, []) : le [] vide = dépendances vides = exécution unique au montage du composant",
    "Sélecteurs ciblés useAppStore(s => s.notes) : Zustand ne re-rend que si notes change, pas tout le store",
]
for i, pt in enumerate(pts2):
    bx = 0.5 + (i % 2) * 6.4
    by = 5.9 + (i // 2) * 0.55
    text_box(slide, "▸  " + pt, bx, by, 6.2, 0.48, size=10, color=GRAY)


# ── SLIDE 27 : Questions fréquentes mentor ────────────────────────────────────
slide = add_slide()
text_box(slide, "Questions que votre mentor peut poser — et les réponses", 0.5, 0.3, 12, 0.7, size=24, bold=True, color=YELLOW)
rect(slide, 0.5, 1.05, 12.3, 0.04, YELLOW)

qas = [
    ("Pourquoi un Service et pas directement dans le Controller ?",
     "Le Controller a pour rôle de gérer HTTP (recevoir/répondre). Si on met la logique dedans, "
     "on viole le S de SOLID. Le Service est réutilisable depuis n'importe où (API, tests, CLI)."),
    ("Quelle est la différence entre MVC et ce qu'on a fait ?",
     "On a fait MVC + Service Layer. Dans un MVC pur, le Controller parle directement au Model. "
     "Ici on ajoute une couche Service entre les deux pour isoler la logique métier."),
    ("Pourquoi Zustand plutôt que Redux ?",
     "Redux exige actions/reducers/dispatch séparés — trop de boilerplate pour un projet de cette taille. "
     "Zustand fait la même chose en 10x moins de code. Le pattern Flux reste le même."),
    ("Comment fonctionne l'authentification dans l'API ?",
     "Le middleware 'auth' de Laravel vérifie la session cookie existante. "
     "Axios envoie le cookie automatiquement via withCredentials:true. "
     "Le token CSRF protège contre les attaques CSRF sur les POST/DELETE."),
    ("Qu'est-ce que l'inversion de dépendances concrètement ?",
     "NoteController ne fait pas new NoteService(). Il déclare en constructeur qu'il a BESOIN "
     "d'un NoteService. Laravel le crée et l'injecte. On dépend de l'abstraction, pas de l'implémentation."),
    ("Quels fichiers Livewire peut-on supprimer maintenant ?",
     "app/Livewire/Notes.php et app/Livewire/TagForm.php sont obsolètes. "
     "Les vues livewire/notes.blade.php et livewire/tag-form.blade.php aussi. "
     "On les garde pour l'instant pour ne pas casser la migration progressive."),
]
for i, (q, a) in enumerate(qas):
    bx = 0.4 + (i % 2) * 6.4
    by = 1.25 + (i // 2) * 1.95
    rect(slide, bx, by, 6.0, 1.8, RGBColor(0x31, 0x32, 0x44))
    rect(slide, bx, by, 6.0, 0.06, YELLOW)
    text_box(slide, "Q : " + q, bx+0.15, by+0.12, 5.7, 0.55, size=10, bold=True, color=YELLOW)
    rect(slide, bx+0.1, by+0.72, 5.8, 0.02, GRAY)
    text_box(slide, "R : " + a,  bx+0.15, by+0.82, 5.7, 0.85, size=10, color=WHITE)


# ── SLIDE 28 : Ce qu'on aurait pu faire en plus ───────────────────────────────
slide = add_slide()
text_box(slide, "Pistes d'amélioration — aller plus loin", 0.5, 0.3, 12, 0.7, size=26, bold=True, color=GRAY)
rect(slide, 0.5, 1.05, 12.3, 0.04, GRAY)
text_box(slide, "Ce qui n'a pas été implémenté dans ce projet mais serait pertinent en production :", 0.5, 1.15, 12, 0.4, size=13, color=GRAY)

improvements = [
    (ACCENT, "Repository Pattern",
     "Ajouter une couche Repository entre Service et Model.\n"
     "NoteService ne parle plus à Note::create() directement\n"
     "mais à NoteRepository::create().\n"
     "→ Encore plus testable, abstrait l'ORM."),
    (GREEN, "Laravel Sanctum (tokens API)",
     "Remplacer le middleware 'auth' session\n"
     "par des tokens Bearer (Sanctum).\n"
     "Permet à une app mobile de s'authentifier\n"
     "sans cookies de session."),
    (YELLOW, "Form Requests Laravel",
     "Déplacer la validation du Controller\n"
     "dans des classes dédiées (StoreNoteRequest).\n"
     "Principe S : le Controller ne s'occupe plus\n"
     "de la validation, juste du routing."),
    (ORANGE, "Tests automatisés",
     "PHPUnit pour tester NoteService sans HTTP.\n"
     "Jest/Vitest pour tester les composants React.\n"
     "Tests d'intégration sur les endpoints API\n"
     "avec Laravel HTTP Testing."),
    (RED, "Documentation Swagger/OpenAPI",
     "Documenter l'API avec des annotations\n"
     "ou un fichier openapi.yaml.\n"
     "Génère une interface interactive\n"
     "pour tester les endpoints."),
    (ACCENT, "Soft Deletes",
     "Ajouter deleted_at aux migrations\n"
     "et SoftDeletes aux Models.\n"
     "Les notes 'supprimées' restent en BDD\n"
     "pour un éventuel historique."),
]
for i, (col, title, desc) in enumerate(improvements):
    bx = 0.4 + (i % 3) * 4.25
    by = 1.7 + (i // 3) * 2.65
    rect(slide, bx, by, 3.95, 2.4, RGBColor(0x31, 0x32, 0x44))
    rect(slide, bx, by, 0.1, 2.4, col)
    text_box(slide, title, bx+0.25, by+0.12, 3.55, 0.45, size=13, bold=True, color=col)
    text_box(slide, desc,  bx+0.25, by+0.65, 3.55, 1.6,  size=11, color=WHITE)


# ── Save ─────────────────────────────────────────────────────────────────────
output = "/home/sacha/Public/openclassroom/architect/P3_architect/analyse_architecture_renote.pptx"
prs.save(output)
print("Saved:", output)
